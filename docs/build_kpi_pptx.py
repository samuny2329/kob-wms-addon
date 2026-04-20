"""
KPI Assessment Guide — PowerPoint Walkthrough (Screenshot Edition)

Each slide = 1 real screenshot + title + caption.
45 screenshots arranged in logical step-by-step order.

Usage:
    odoo-19.0/venv/Scripts/pip install python-pptx Pillow
    odoo-19.0/venv/Scripts/python custom_addons/kob_wms/docs/build_kpi_pptx.py

Output:
    docs/KPI_Assessment_Guide.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

try:
    from PIL import Image as PILImage
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# ── Paths ──────────────────────────────────────────────────────────────────
DOCS_DIR  = Path(__file__).parent
SHOT_DIR  = DOCS_DIR.parent / "kpi_screenshots"
PPTX_PATH = DOCS_DIR / "KPI_Assessment_Guide.pptx"

# ── Slide dimensions: Widescreen 13.33" × 7.5" ────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

# ── Colours ────────────────────────────────────────────────────────────────
KOB_PURPLE  = RGBColor(0x4B, 0x00, 0x82)   # main purple header
KOB_DARK    = RGBColor(0x2C, 0x00, 0x5C)   # caption bar
KOB_LIGHT   = RGBColor(0x7B, 0x2F, 0xBE)   # accent / title slide gradient
GOLD        = RGBColor(0xFF, 0xD0, 0x00)   # section label accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG    = RGBColor(0xF5, 0xF0, 0xFF)   # slide background
GRAY_TEXT   = RGBColor(0x44, 0x44, 0x44)

# ── Layout constants ───────────────────────────────────────────────────────
HEADER_H  = Inches(0.60)   # purple top bar
CAPTION_H = Inches(0.85)   # dark caption bar at bottom
IMG_TOP   = HEADER_H
IMG_H     = H - HEADER_H - CAPTION_H   # ≈ 6.05"
IMG_W     = W                           # full width

# ─────────────────────────────────────────────────────────────────────────
# SLIDES: (filename_stem, title, caption)
# Ordered for logical step-by-step flow through the system
# ─────────────────────────────────────────────────────────────────────────
SLIDES = [

    # ═══ SECTION 1: System Dashboards ════════════════════════════════════
    (
        "Screenshot 2026-04-20 010100",
        "Dashboard — WMS Overview",
        "ภาพรวม KOB WMS: Orders by Status (bar), Orders by Platform (pie), Daily Order Volume (line), Orders by Courier "
        "และ Pick/Pack Errors by Status — ดู real-time ได้จากเมนู Dashboard",
    ),
    (
        "Screenshot 2026-04-20 010114",
        "Dashboard — KPI Performance",
        "Worker Score Top Performers, Actions by Worker, UPH Trend รายวัน, Errors by Worker และ Quality % by Worker "
        "— ใช้ติดตามประสิทธิภาพพนักงานแบบ real-time กรองตาม Period และ Worker ได้",
    ),
    (
        "Screenshot 2026-04-20 010131",
        "Dashboard — Inventory Count Adjustments",
        "Adjustments by State (pie), Total Variance by Product, Adjustments by Location, Count Sessions by Type "
        "และ Daily Adjustments Created — ติดตามผลการนับสต็อก Cycle Count",
    ),
    (
        "Screenshot 2026-04-20 010142",
        "Dashboard — WMS Operations",
        "Avg Pick Duration by Picker (min), Avg Pack Duration by Packer (min), Count Sessions by State, "
        "Pickfaces by Zone และ Daily Shipped Orders — วัดเวลา Pick/Pack ต่อคนและติดตาม shipments",
    ),
    (
        "Screenshot 2026-04-19 185451",
        "WMS Overview Dashboard (Full View)",
        "Dashboard Overview แสดงใน Odoo พร้อม sidebar KOB WMS: WMS Overview, KPI Performance, "
        "Count Adjustments, WMS Operations — เข้าถึงได้จากเมนู Dashboard ในแถบบนสุด",
    ),
    (
        "Screenshot 2026-04-19 185708",
        "WMS Overview — Browser Context",
        "หน้า Dashboard เดียวกัน พร้อม Chrome Extensions panel — ระบบรองรับการใช้งานบน Chrome browser ปกติ "
        "ไม่ต้องติดตั้ง software พิเศษ",
    ),
    (
        "Screenshot 2026-04-19 190333",
        "WMS Overview (Compact Window)",
        "Dashboard แสดงในหน้าต่างขนาดเล็กพร้อม Assistant sidebar — ระบบ responsive "
        "และใช้งานได้หลายหน้าต่างพร้อมกัน",
    ),
    (
        "Screenshot 2026-04-19 191855",
        "KPI Performance Dashboard (Sidebar View)",
        "KPI Performance Dashboard: Worker Score, Actions, UPH Trend, Errors, Quality % "
        "— แสดงพร้อม assistant sidebar สำหรับ consultation และ analysis",
    ),

    # ═══ SECTION 2: Fulfilment Workflow ══════════════════════════════════
    (
        "Screenshot 2026-04-19 200337",
        "Fulfilment Step 1 — Pick Queue (F1)",
        "รายการคำสั่งซื้อที่รอ Pick: Order, Platform Ref, Customer, Platform badge (Shopee/Lazada/TikTok/Manual), "
        "Courier, AWB, Box Barcode, Picker และ SLA Status (On Track / At Risk / Breached)",
    ),
    (
        "Screenshot 2026-04-19 200444",
        "Fulfilment Step 1 — Sales Order Detail (Picking)",
        "SO/2026/00596: สแกน SKU barcode เพื่อ Pick สินค้า — แสดง Items tab (Expected/Picked/Packed qty), "
        "Smart Ring SLA tab บันทึก SLA timeline ทุก step และ chatter ทางขวา",
    ),
    (
        "Screenshot 2026-04-19 200717",
        "Fulfilment Step 2 — Pack Queue (F2)",
        "Pack Queue: SO/2026/00592 พร้อม Pack — Platform: Manual, Picker: Picker 01, Items: 2, "
        "Picked: 2, SLA: On Track, Status: Picked — สแกน order barcode เพื่อเปิด",
    ),
    (
        "Screenshot 2026-04-19 200849",
        "Fulfilment Step 2 — Order in Picked State",
        "SO/2026/00592 state Picked: Items (KTAP088 Eau de Toilette + KLAP200 Perfume Lotion qty 1 each), "
        "chatter บันทึก Pending → Picking → Picked พร้อม timestamp ทุก transition",
    ),
    (
        "Screenshot 2026-04-19 201044",
        "Fulfilment Step 2 — Close Box (Box Recommender)",
        "Close Box Wizard: ระบบแนะนำ Box B — 17×25×9 cm (BY COUNT เพราะไม่มี dimension data) "
        "พร้อม Override dropdown สำหรับเลือก box ขนาดอื่น และปุ่ม 'Close Box & Print AWB'",
    ),
    (
        "Screenshot 2026-04-19 201308",
        "Fulfilment Step 3 — All Orders (SLA View)",
        "All Orders list: Order No., Platform Ref, Customer, Courier, AWB/Tracking, Box Barcode, Packer, "
        "SLA Status badge (เขียว=On Track, เหลือง=At Risk, แดง=Breached) และ Status badge",
    ),
    (
        "Screenshot 2026-04-19 201448",
        "Fulfilment — Order Timeline Log",
        "ตาราง orders พร้อม Timestamp, Platform Ref, Order No., Items count, Code Scanned และ Sale Order link "
        "— ใช้ตรวจสอบ SLA compliance และ audit trail ย้อนหลัง",
    ),

    # ═══ SECTION 3: Analytics ════════════════════════════════════════════
    (
        "Screenshot 2026-04-19 201611",
        "Analytics — Worker Performance Report",
        "ตาราง Worker Performance: Date, Employee, Picks, Packs, Boxes, Ships, Total Actions, UPH (8h), "
        "Pick Errors, Pack Errors, Total Error %, Quality % และ Worker Score — export ได้",
    ),
    (
        "Screenshot 2026-04-19 201818",
        "Analytics — Product in Box Analytics",
        "Box Analytics ระดับ product: Product, Box Used, Qty, Lot/Serial, Avg Fill%, Avg Box Time, "
        "Rubber Band, Tape Cost, Bubble Cost, Total Material Cost, Total Pack Cost — วิเคราะห์ต้นทุนบรรจุภัณฑ์",
    ),
    (
        "Screenshot 2026-04-19 202006",
        "Analytics — Box Catalogue (28 THE BOX Sizes)",
        "ตาราง Box ทั้ง 28 ขนาด: Box Photo, Code, Display Label, L/W/H (cm), Volume (cm³ และ m³), "
        "Unit Cost, Tape Cost Est., Bubble Cost Est., Total Material Cost และ In Stock flag",
    ),

    # ═══ SECTION 4: KPI Configuration (Admin Setup) ══════════════════════
    (
        "Screenshot 2026-04-20 010158",
        "KPI Config Step 1 — Seasons List",
        "KPI Seasons: Q1 2026 (Closed, 04/11–04/12, 0 Assessments) และ Q1 2026#01 (Open, 04/11–04/18, 20 Assessments) "
        "— ไปที่ KPI Assessment > Configuration > Seasons เพื่อสร้าง season ใหม่",
    ),
    (
        "Screenshot 2026-04-20 010220",
        "KPI Config Step 1 — Season Form",
        "Season Q1 2026#01: Date Start 04/11, Date End 04/18, Season Type: Half Year (H1/H2), "
        "Assessment Count: 20, Self Weight 40%, Reviewer Weight 60%, Company: บริษัท คีสออฟบิวตี้ จำกัด",
    ),
    (
        "Screenshot 2026-04-20 010238",
        "KPI Config Step 1 — Season Type Options",
        "Season Type dropdown: Full Year, Half Year (H1/H2) ✓ (เลือกอยู่), Mid-Year Check-in "
        "— เลือก Type ตามรอบการประเมินของบริษัท แล้วกด 'Create All Assessments' เพื่อสร้างให้พนักงานทุกคน",
    ),
    (
        "Screenshot 2026-04-20 010316",
        "KPI Config Step 2 — Templates List",
        "KPI Templates: 9 ตำแหน่ง (Admin Online, Driver/Transport, Inbound/Receiving, Inventory/Cycle Count, "
        "Manager, Packer, Picker, Shipper/Outbound, Supervisor) — ทุกตำแหน่ง Total Weight = 100.00%",
    ),
    (
        "Screenshot 2026-04-19 195837",
        "KPI Config Step 2 — Templates Overview",
        "KPI Templates list ครบ 9 ตำแหน่ง พร้อม Total Weight 100.00 ทุกตำแหน่ง "
        "— แต่ละ Template กำหนด Pillar Weights และ SOP Documents ต่างกันตามบทบาทงาน",
    ),
    (
        "Screenshot 2026-04-19 195918",
        "KPI Config — Inbound Template Pillar Weights",
        "Inbound/Receiving Template: 8 Pillars พร้อม Weight% — Develop Our People 30%, Drive Value 15%, "
        "Make Revenue 10%, Champion Progress 15%, Deliver Financial 5%, Manage Risk 10%, Live Values 5%, Ops Excellence 10%",
    ),
    (
        "Screenshot 2026-04-20 010350",
        "KPI Config — Supervisor Template SOPs",
        "Supervisor Template > SOP Documents tab: SOP-SV-001 'การบริหารทีมประจำกะ' Version 1.0, "
        "Related Pillar: Develop Our People, Active ✓ — SOP อ้างอิงสำหรับพนักงานขณะประเมิน",
    ),
    (
        "Screenshot 2026-04-20 010419",
        "KPI Config — Packer Template SOPs",
        "Packer Template > SOP Documents: SOP-PA-001 'ขั้นตอนการ Pack สินค้า' v2.0 (Develop Our People) "
        "และ SOP-PA-002 'มาตรฐาน QC ก่อนปิดกล่อง' v1.0 (Drive Value for Clients) — Active ทั้งคู่",
    ),
    (
        "Screenshot 2026-04-20 010516",
        "KPI Config — Pillar Weights & Criteria Dialog",
        "Supervisor Template > Pillar Weights tab: 8 Pillars listed. Dialog 'Open: Pillar Weights' — "
        "Develop Our People 15%: Criteria = Team Development & Training 50%, Performance Management 50%",
    ),
    (
        "Screenshot 2026-04-19 200104",
        "KPI Config Step 3 — Approver Configuration",
        "Approver Config: 9 ตำแหน่ง กำหนด Default Supervisor, Asst.Manager, Manager, Director ประจำแต่ละ Position "
        "— ค่า default นี้จะถูก assign ให้ Assessment อัตโนมัติเมื่อ 'Create All Assessments'",
    ),

    # ═══ SECTION 5: Assessment Process (Employee Journey) ════════════════
    (
        "Screenshot 2026-04-19 192217",
        "Assessment Step 1 — Draft Created",
        "My Assessment: Admin KOB, Season Q1 2026#01, Position: Admin Online, Grade E (ยังไม่มีคะแนน), "
        "State: Draft — พนักงานเห็น assessment ของตัวเองหลัง Supervisor กด 'Create All Assessments'",
    ),
    (
        "Screenshot 2026-04-19 192231",
        "Assessment Step 1 — Assessment Form Overview",
        "Assessment Form: Employee, Position, Season, Approvers (Supervisor / Asst.Manager / Manager / Director) "
        "— Tabs: Pillar Scores, Comments, SOP References, Goals & Action Items",
    ),
    (
        "Screenshot 2026-04-19 192316",
        "Assessment Step 2 — Pillar Scores Tab",
        "Pillar Scores: 8 Pillars แต่ละแถวมี Weight%, Self Score (พนักงานกรอก), "
        "Reviewer Score (Supervisor กรอก), Weighted Score (คำนวณอัตโนมัติ) — คลิกแถวเพื่อเปิด dialog",
    ),
    (
        "Screenshot 2026-04-19 192346",
        "Assessment Step 2 — Criterion Score Dialog",
        "Dialog Pillar 'Develop Our People': รายการ Criteria (Talent Development, Team Motivation) "
        "— แต่ละ Criterion มี Self radio (0–5) และ Reviewer radio (0–5) พร้อมคลิก Rubric เพื่อดูคำอธิบาย",
    ),
    (
        "Screenshot 2026-04-19 192406",
        "Assessment Step 2 — Scoring Rubric (1–5 Scale)",
        "Rubric สำหรับ Talent Development: Score 1 (ยังไม่ดำเนินการ) → Score 3 (ดำเนินการสม่ำเสมอ) "
        "→ Score 5 (เป็น mentor ระดับองค์กร) — พนักงานเลือก Self radio, Supervisor เลือก Reviewer radio",
    ),
    (
        "Screenshot 2026-04-19 192452",
        "Assessment Step 3 — Comments Tab",
        "Comments Tab: 5 ช่องความคิดเห็นแยกตามระดับ — Employee, Supervisor, Asst.Manager, Manager, Director "
        "— บันทึก feedback ระหว่างขั้นตอน approval chain",
    ),
    (
        "Screenshot 2026-04-19 192520",
        "Assessment Step 4 — SOP Reference Popup",
        "SOP Reference dialog: SOP-AO-001 ขั้นตอน Admin Online — แสดง Platform SLA Table "
        "(Shopee/Lazada/TikTok cut-off times) — พนักงานอ่านก่อนประเมินเพื่อให้คะแนนถูกต้อง",
    ),
    (
        "Screenshot 2026-04-19 192534",
        "Assessment Step 4 — SOP Document (Full View)",
        "SOP-AO-001 form เต็ม: Title, Version 2.0, Effective Date, Related Pillar: Drive Value for Clients, "
        "Content (Platform SLA table พร้อม cut-off times ต่อ platform) — Active ✓",
    ),
    (
        "Screenshot 2026-04-19 192608",
        "Assessment Step 5 — Goals & Action Items",
        "Goals tab: 5 Goals แต่ละรายการมี Related Pillar, Goal Description, Deadline, State (To Do) "
        "— พนักงานตั้งเป้าหมายและ action items สำหรับ period นี้เพื่อ improvement",
    ),
    (
        "Screenshot 2026-04-19 195632",
        "Assessment — Criterion Score with Comments",
        "Innovation & Automation pillar: Self Score = 5, Reviewer Score = 4 "
        "พร้อม Thai comments จาก Supervisor — บันทึกเหตุผลและ feedback การให้คะแนน",
    ),

    # ═══ SECTION 6: Approval Chain & Completion ══════════════════════════
    (
        "Screenshot 2026-04-19 195115",
        "Assessment Completed — Grade A (Final 4.51)",
        "My Assessment: Admin KOB, Season Q1 2026#01, Self Score 5.00, Supervisor Score 4.18, "
        "Final Score 4.51 (Self×40% + Supervisor×60%), Grade A (≥4.5), State: Completed",
    ),
    (
        "Screenshot 2026-04-19 195135",
        "Assessment Completed — Approval Chain in Chatter",
        "Chatter บันทึก approval chain ครบ: Draft → Self Assessment → Supervisor Review "
        "→ Asst.Manager Review → Manager Approval → Director Approval → Completed พร้อม timestamp ทุก step",
    ),
    (
        "Screenshot 2026-04-19 195809",
        "Completed Assessment — SOP References Tab",
        "SOP References tab ของ Assessment ที่ Completed: แสดง SOPs ที่ใช้อ้างอิงระหว่างการประเมิน "
        "พร้อม activity notifications ใน chatter ด้านขวาแสดง approval history",
    ),
    (
        "Screenshot 2026-04-19 195820",
        "Completed Assessment — Goals Tab",
        "Goals & Action Items tab ของ Assessment ที่ Completed: แสดง goals ที่ตั้งไว้พร้อมสถานะ "
        "— chatter ด้านขวาบันทึกประวัติ approval ครบทุกขั้นตั้งแต่ Draft ถึง Completed",
    ),
    (
        "Screenshot 2026-04-19 200002",
        "All Assessments — Season Overview (20 Records)",
        "All Assessments รวม 20 รายการ: 1 Completed (Admin KOB Final 4.51 Grade A) + 19 Draft "
        "(Pickers, Packers, Supervisors) — Manager เห็นภาพรวม assessment ทุกคนในบริษัท",
    ),

    # ═══ SECTION 7: History & Calibration ══════════════════════════════
    (
        "Screenshot 2026-04-20 005927",
        "Assessment History",
        "Assessment History: Admin KOB, Season Q1 2026#01, Position: Admin Online, "
        "Self Score 5.00, Supervisor Score 4.18, Final 4.51, Grade A, State: Completed "
        "— ดูประวัติการประเมินทุก Season ย้อนหลังได้ที่ KPI Assessment > Assessment History",
    ),
    (
        "Screenshot 2026-04-20 010016",
        "KPI Calibration — Pivot View",
        "KPI Calibration Pivot: Season Q1 2026#01 → Position Admin Online: Final 4.51, Self 5.00, Supervisor 4.18 "
        "— Managers ใช้ Calibration Pivot เพื่อ compare scores ข้าม position และปรับให้สม่ำเสมอทั้งองค์กร",
    ),
]


# ─────────────────────────────────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────────────────────────────────

def _blank(prs):
    """Add blank slide (layout 6 = completely blank)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, width, height, rgb):
    """Add a filled rectangle with no border."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    s.line.fill.background()
    return s


def _tb(slide, left, top, width, height, text, size,
        bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def _fit_image(img_path, max_w_emu, max_h_emu):
    """Return (width_emu, height_emu) that fits within max bounds, preserving aspect ratio."""
    if HAS_PIL:
        with PILImage.open(img_path) as im:
            iw, ih = im.size
    else:
        # Fallback: assume 16:9
        iw, ih = 1920, 1080

    scale = min(max_w_emu / iw, max_h_emu / ih)
    return int(iw * scale), int(ih * scale)


# ─────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────

def slide_title(prs):
    """Cover slide."""
    s = _blank(prs)

    # Background
    _rect(s, 0, 0, W, H, KOB_PURPLE)

    # Decorative accent strip at top
    _rect(s, 0, 0, W, Inches(0.12), GOLD)

    # Decorative stripe at bottom
    _rect(s, 0, H - Inches(0.12), W, Inches(0.12), GOLD)

    # Company logo area — white box on left
    _rect(s, Inches(0.5), Inches(1.6), Inches(3.5), Inches(1.0), WHITE)
    _tb(s, Inches(0.5), Inches(1.65), Inches(3.5), Inches(0.9),
        "KOB WMS Pro", 22, bold=True, color=KOB_PURPLE, align=PP_ALIGN.CENTER)

    # Main title
    _tb(s, Inches(0.4), Inches(2.9), Inches(12.5), Inches(1.4),
        "KPI Assessment Guide", 52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    _tb(s, Inches(0.4), Inches(4.4), Inches(12.5), Inches(0.7),
        "ขั้นตอนการประเมิน KPI — Kiss of Beauty (KOB) / SKINOXY", 22,
        bold=False, color=GOLD, align=PP_ALIGN.LEFT)

    # Description
    _tb(s, Inches(0.4), Inches(5.2), Inches(9.0), Inches(0.7),
        "45 ภาพหน้าจอจริงจากระบบ · เรียงตามลำดับขั้นตอน · Odoo 18 Community",
        16, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    # Section overview (right side)
    sections = [
        "① System Dashboards (8 slides)",
        "② Fulfilment Workflow (7 slides)",
        "③ Analytics (3 slides)",
        "④ KPI Configuration (9 slides)",
        "⑤ Assessment Process (10 slides)",
        "⑥ Approval & Completion (5 slides)",
        "⑦ History & Calibration (2 slides)",
        "⑧ Assessment History (1 slide)",
    ]
    _tb(s, Inches(9.8), Inches(1.5), Inches(3.2), Inches(5.5),
        "\n".join(sections), 11, bold=False, color=WHITE, align=PP_ALIGN.LEFT)


def slide_screenshot(prs, img_path, step_num, total, title, caption):
    """One screenshot slide with header, image, and caption bar."""
    s = _blank(prs)

    # ── Background ─────────────────────────────────────────────────────
    _rect(s, 0, 0, W, H, LIGHT_BG)

    # ── Header bar ─────────────────────────────────────────────────────
    _rect(s, 0, 0, W, HEADER_H, KOB_PURPLE)

    # Step counter pill (gold)
    pill_w = Inches(1.1)
    _rect(s, Inches(0.12), Inches(0.10), pill_w, Inches(0.40), GOLD)
    _tb(s, Inches(0.12), Inches(0.10), pill_w, Inches(0.40),
        f"STEP {step_num:02d}/{total}", 9, bold=True, color=KOB_PURPLE,
        align=PP_ALIGN.CENTER)

    # Title in header
    _tb(s, Inches(1.35), Inches(0.07), Inches(11.8), Inches(0.47),
        title, 15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # ── Image ──────────────────────────────────────────────────────────
    max_w = int(IMG_W)
    max_h = int(IMG_H)
    try:
        img_w, img_h = _fit_image(img_path, max_w, max_h)
        # Center image horizontally and vertically within the image area
        left = (max_w - img_w) // 2
        top  = int(IMG_TOP) + (max_h - img_h) // 2
        s.shapes.add_picture(str(img_path), left, top, img_w, img_h)
    except Exception as e:
        # Fallback: show error text
        _tb(s, Inches(0.5), Inches(2.0), Inches(12.0), Inches(1.0),
            f"[Image not found: {img_path.name}]  Error: {e}", 14,
            bold=False, color=RGBColor(0xCC, 0x00, 0x00))

    # ── Caption bar ────────────────────────────────────────────────────
    cap_top = H - CAPTION_H
    _rect(s, 0, cap_top, W, CAPTION_H, KOB_DARK)

    _tb(s, Inches(0.25), cap_top + Inches(0.08), Inches(12.8), CAPTION_H - Inches(0.10),
        caption, 12, bold=False, color=WHITE, align=PP_ALIGN.LEFT)


def slide_closing(prs):
    """Thank-you / summary closing slide."""
    s = _blank(prs)

    _rect(s, 0, 0, W, H, KOB_PURPLE)
    _rect(s, 0, 0, W, Inches(0.12), GOLD)
    _rect(s, 0, H - Inches(0.12), W, Inches(0.12), GOLD)

    _tb(s, Inches(0.5), Inches(1.8), Inches(12.0), Inches(1.4),
        "KOB WMS Pro — KPI Assessment", 44, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

    _tb(s, Inches(0.5), Inches(3.4), Inches(12.0), Inches(0.8),
        "ขั้นตอนครบวงจร: Configuration → Assessment → Approval → Calibration",
        22, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

    summary = (
        "Grade Scale:  A (≥4.5)  ·  B (≥3.5)  ·  C (≥2.5)  ·  D (≥1.5)  ·  E (<1.5)\n"
        "Final Score = Self × 40%  +  Supervisor × 60%\n"
        "Approval Chain: Draft → Self Assessment → Supervisor → Asst.Mgr → Manager → Director → Completed"
    )
    _tb(s, Inches(0.5), Inches(4.4), Inches(12.0), Inches(1.5),
        summary, 15, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    _tb(s, Inches(0.5), Inches(6.1), Inches(12.0), Inches(0.6),
        "Kiss of Beauty (KOB) / SKINOXY  ·  Odoo 18 Community  ·  kob_wms v18.0.2.15.0",
        13, bold=False, color=RGBColor(0xCC, 0xBB, 0xFF), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Cover
    slide_title(prs)
    print("✓ Title slide")

    # Screenshot slides
    total = len(SLIDES)
    missing = []

    for i, (stem, title, caption) in enumerate(SLIDES, start=1):
        img_path = SHOT_DIR / f"{stem}.png"
        if not img_path.exists():
            print(f"  ✗ MISSING: {img_path.name}")
            missing.append(img_path.name)
        else:
            print(f"  [{i:02d}/{total}] {title}")
        slide_screenshot(prs, img_path, i, total, title, caption)

    # Closing
    slide_closing(prs)
    print("✓ Closing slide")

    # Save
    prs.save(str(PPTX_PATH))
    print(f"\n✅ Saved → {PPTX_PATH}")
    print(f"   Total slides: {total + 2}  (1 title + {total} screenshots + 1 closing)")

    if missing:
        print(f"\n⚠️  {len(missing)} missing images:")
        for m in missing:
            print(f"   - {m}")


if __name__ == "__main__":
    main()
